"""
PowerPoint Generator for Signal
Creates TAU-branded presentations from data/prompts
Enhanced with charts, gradient backgrounds, and professional layouts.
"""
import os
import io
import json
import base64
import tempfile
from datetime import datetime
from pathlib import Path
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional, List, Dict, Any
import anthropic

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
    from lxml import etree
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import matplotlib
    matplotlib.use("Agg")  # headless backend
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    import numpy as np
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False

router = APIRouter(prefix="/api", tags=["ppt"])

# ─── TAU Colour Palette ───────────────────────────────────────────────────────
TAU_NAVY    = RGBColor(0x00, 0x33, 0x66) if PPTX_AVAILABLE else None
TAU_BLUE    = RGBColor(0x33, 0x85, 0xD6) if PPTX_AVAILABLE else None
TAU_LIGHT   = RGBColor(0x66, 0xA3, 0xE0) if PPTX_AVAILABLE else None
TAU_WHITE   = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
TAU_DARK    = RGBColor(0x1A, 0x1A, 0x2E) if PPTX_AVAILABLE else None
TAU_ACCENT  = RGBColor(0x00, 0xC4, 0xCC) if PPTX_AVAILABLE else None   # teal highlight
TAU_GOLD    = RGBColor(0xFF, 0xC1, 0x07) if PPTX_AVAILABLE else None   # chart accent
TAU_GREY    = RGBColor(0x8A, 0x9B, 0xAB) if PPTX_AVAILABLE else None

# Hex versions for matplotlib
CHART_PALETTE = ["#3385D6", "#00C4CC", "#FFC107", "#003366", "#66A3E0", "#E84393"]

ANTHROPIC_KEY = os.getenv("ANTHROPIC_API_KEY")

# ─── Models ──────────────────────────────────────────────────────────────────

class ChartData(BaseModel):
    chart_type: str = "bar"          # bar | line | pie | grouped_bar
    title: str = ""
    labels: List[str] = []
    datasets: List[Dict[str, Any]] = []   # [{"label": "...", "data": [...]}]
    x_label: Optional[str] = None
    y_label: Optional[str] = None

class SlideContent(BaseModel):
    title: str
    bullets: List[str] = []
    notes: Optional[str] = None
    layout: str = "content"          # title | content | chart | closing | section
    chart_data: Optional[ChartData] = None

class RetailMetrics(BaseModel):
    """Convenience model for retail media KPIs."""
    channels: List[str] = ["Sponsored Products", "Display", "Video", "DSP"]
    spend:       List[float] = [125000, 45000, 30000, 20000]
    revenue:     List[float] = [875000, 198000, 90000, 56000]
    roas:        List[float] = [7.0, 4.4, 3.0, 2.8]
    impressions: List[float] = [12500000, 8200000, 3100000, 5400000]
    # Month-over-month trend (last 6 months)
    trend_months: List[str] = ["Oct", "Nov", "Dec", "Jan", "Feb", "Mar"]
    trend_spend:  List[float] = [160000, 175000, 210000, 190000, 205000, 220000]
    trend_revenue:List[float] = [960000, 1050000, 1470000, 1140000, 1230000, 1540000]

class PPTRequest(BaseModel):
    title: str
    subtitle: Optional[str] = None
    company: Optional[str] = None
    prompt: Optional[str] = None
    slides: Optional[List[SlideContent]] = None
    context_data: Optional[dict] = None
    include_signal_data: bool = False
    retail_metrics: Optional[RetailMetrics] = None


# ─── Chart helpers ────────────────────────────────────────────────────────────

def _fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def make_bar_chart(title: str, labels: List[str], values: List[float],
                   ylabel: str = "", color: str = "#3385D6") -> bytes:
    fig, ax = plt.subplots(figsize=(9, 5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    x = np.arange(len(labels))
    bars = ax.bar(x, values, color=CHART_PALETTE[:len(labels)], edgecolor="none",
                  width=0.6, zorder=3)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, color="white", fontsize=11)
    ax.set_ylabel(ylabel, color="#66A3E0", fontsize=11)
    ax.set_title(title, color="white", fontsize=14, fontweight="bold", pad=12)
    ax.tick_params(colors="#66A3E0")
    ax.yaxis.grid(True, color="#1E3A5F", linestyle="--", alpha=0.6, zorder=0)
    ax.set_axisbelow(True)
    for spine in ax.spines.values():
        spine.set_visible(False)
    # Value labels on bars
    for bar in bars:
        h = bar.get_height()
        label = f"{h:,.0f}" if h >= 1000 else f"{h:.2f}"
        ax.text(bar.get_x() + bar.get_width() / 2, h * 1.02, label,
                ha="center", va="bottom", color="white", fontsize=10, fontweight="bold")
    fig.tight_layout()
    return _fig_to_bytes(fig)


def make_grouped_bar_chart(title: str, labels: List[str],
                            datasets: List[Dict], ylabel: str = "") -> bytes:
    n_groups = len(labels)
    n_series = len(datasets)
    width = 0.7 / n_series
    fig, ax = plt.subplots(figsize=(10, 5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    x = np.arange(n_groups)
    for i, ds in enumerate(datasets):
        offset = (i - n_series / 2 + 0.5) * width
        bars = ax.bar(x + offset, ds["data"], width=width,
                      color=CHART_PALETTE[i % len(CHART_PALETTE)],
                      label=ds.get("label", f"Series {i+1}"), edgecolor="none", zorder=3)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, color="white", fontsize=10)
    ax.set_ylabel(ylabel, color="#66A3E0", fontsize=11)
    ax.set_title(title, color="white", fontsize=14, fontweight="bold", pad=12)
    ax.tick_params(colors="#66A3E0")
    ax.yaxis.grid(True, color="#1E3A5F", linestyle="--", alpha=0.6, zorder=0)
    ax.set_axisbelow(True)
    for spine in ax.spines.values():
        spine.set_visible(False)
    legend = ax.legend(facecolor="#0D1B2A", labelcolor="white", fontsize=10,
                       framealpha=0.8, edgecolor="#1E3A5F")
    fig.tight_layout()
    return _fig_to_bytes(fig)


def make_line_chart(title: str, x_labels: List[str],
                    datasets: List[Dict], ylabel: str = "") -> bytes:
    fig, ax = plt.subplots(figsize=(9, 5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    for i, ds in enumerate(datasets):
        color = CHART_PALETTE[i % len(CHART_PALETTE)]
        ax.plot(x_labels, ds["data"], marker="o", color=color,
                label=ds.get("label", f"Series {i+1}"),
                linewidth=2.5, markersize=7)
        ax.fill_between(x_labels, ds["data"], alpha=0.12, color=color)
    ax.set_ylabel(ylabel, color="#66A3E0", fontsize=11)
    ax.set_title(title, color="white", fontsize=14, fontweight="bold", pad=12)
    ax.tick_params(colors="#66A3E0")
    ax.yaxis.grid(True, color="#1E3A5F", linestyle="--", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)
    for spine in ax.spines.values():
        spine.set_visible(False)
    plt.setp(ax.get_xticklabels(), color="white", fontsize=11)
    if len(datasets) > 1:
        legend = ax.legend(facecolor="#0D1B2A", labelcolor="white", fontsize=10,
                           framealpha=0.8, edgecolor="#1E3A5F")
    fig.tight_layout()
    return _fig_to_bytes(fig)


def make_pie_chart(title: str, labels: List[str], values: List[float]) -> bytes:
    fig, ax = plt.subplots(figsize=(7, 5), facecolor="#0D1B2A")
    ax.set_facecolor("#0D1B2A")
    wedges, texts, autotexts = ax.pie(
        values, labels=None,
        colors=CHART_PALETTE[:len(labels)],
        autopct="%1.1f%%", startangle=140,
        pctdistance=0.75,
        wedgeprops={"edgecolor": "#0D1B2A", "linewidth": 2}
    )
    for t in autotexts:
        t.set_color("white")
        t.set_fontsize(10)
    ax.set_title(title, color="white", fontsize=14, fontweight="bold", pad=12)
    patches = [mpatches.Patch(color=CHART_PALETTE[i], label=labels[i])
               for i in range(len(labels))]
    ax.legend(handles=patches, loc="lower center", bbox_to_anchor=(0.5, -0.12),
              facecolor="#0D1B2A", labelcolor="white", fontsize=9,
              ncol=min(len(labels), 4), framealpha=0.8, edgecolor="#1E3A5F")
    fig.tight_layout()
    return _fig_to_bytes(fig)


def chart_bytes_from_spec(cd: ChartData) -> Optional[bytes]:
    if not MATPLOTLIB_AVAILABLE:
        return None
    try:
        if cd.chart_type == "pie":
            vals = cd.datasets[0]["data"] if cd.datasets else []
            return make_pie_chart(cd.title, cd.labels, vals)
        elif cd.chart_type == "line":
            return make_line_chart(cd.title, cd.labels, cd.datasets,
                                   ylabel=cd.y_label or "")
        elif cd.chart_type == "grouped_bar":
            return make_grouped_bar_chart(cd.title, cd.labels, cd.datasets,
                                          ylabel=cd.y_label or "")
        else:  # bar (default)
            vals = cd.datasets[0]["data"] if cd.datasets else []
            return make_bar_chart(cd.title, cd.labels, vals,
                                  ylabel=cd.y_label or "")
    except Exception as e:
        print(f"Chart render error: {e}")
        return None


# ─── Background / design helpers ─────────────────────────────────────────────

def _hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def add_solid_background(slide, rgb: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_gradient_background(slide, from_rgb: tuple, to_rgb: tuple, angle: int = 135):
    """
    Apply a two-stop linear gradient to a slide background via direct XML editing.
    from_rgb / to_rgb: (R, G, B) tuples (0-255).
    """
    r1, g1, b1 = from_rgb
    r2, g2, b2 = to_rgb

    # angle in OOXML is in 60000ths of a degree
    ang = (angle % 360) * 60000

    bg_xml = (
        f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f'      xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'  <p:bgPr>'
        f'    <a:gradFill rot="1">'
        f'      <a:gsLst>'
        f'        <a:gs pos="0">'
        f'          <a:srgbClr val="{r1:02X}{g1:02X}{b1:02X}"/>'
        f'        </a:gs>'
        f'        <a:gs pos="100000">'
        f'          <a:srgbClr val="{r2:02X}{g2:02X}{b2:02X}"/>'
        f'        </a:gs>'
        f'      </a:gsLst>'
        f'      <a:lin ang="{ang}" scaled="0"/>'
        f'    </a:gradFill>'
        f'    <a:effectLst/>'
        f'  </p:bgPr>'
        f'</p:bg>'
    )
    try:
        bg_elem = parse_xml(bg_xml)
        slide._element.insert(0, bg_elem)
    except Exception:
        # Fallback to solid colour
        add_solid_background(slide, RGBColor(r1, g1, b1))


def add_bottom_bar(slide, prs_height, label: str = "TAU — The Independent Marketing Intelligence Architect"):
    """Add a dark footer bar at the bottom of a slide."""
    bar_h = Inches(0.45)
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), prs_height - bar_h,
        prs_height * 2,  # wide enough (width = slide width via prs)
        bar_h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x00, 0x1A, 0x33)
    bar.line.color.rgb = RGBColor(0x00, 0x1A, 0x33)

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0xA3, 0xE0)
    p.alignment = PP_ALIGN.CENTER


def add_accent_line(slide, y_pos, width_inches: float = 12.333,
                    color: RGBColor = None, thickness_pt: float = 3):
    """Horizontal accent rule."""
    color = color or TAU_ACCENT
    line = slide.shapes.add_shape(1,
        Inches(0.5), y_pos,
        Inches(width_inches), Pt(thickness_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.color.rgb = color


def add_text_box(slide, text: str, left, top, width, height,
                 size: int = 20, bold: bool = False,
                 color: RGBColor = None, align=PP_ALIGN.LEFT,
                 italic: bool = False):
    color = color or TAU_WHITE
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def embed_chart_image(slide, img_bytes: bytes, left, top, width, height):
    """Embed a PNG chart image on the slide."""
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, left, top, width, height)


# ─── Retail metrics → chart slides ───────────────────────────────────────────

def build_retail_chart_slides(metrics: RetailMetrics) -> List[SlideContent]:
    slides = []

    # 1. Spend by Channel (bar)
    slides.append(SlideContent(
        title="Spend by Channel",
        layout="chart",
        bullets=["Total spend distributed across key retail media channels"],
        chart_data=ChartData(
            chart_type="bar",
            title="Media Spend by Channel (£)",
            labels=metrics.channels,
            datasets=[{"label": "Spend", "data": metrics.spend}],
            y_label="Spend (£)"
        )
    ))

    # 2. ROAS by Channel (bar)
    slides.append(SlideContent(
        title="ROAS Performance by Channel",
        layout="chart",
        bullets=["Sponsored Products delivers highest return on ad spend"],
        chart_data=ChartData(
            chart_type="bar",
            title="Return on Ad Spend (ROAS) by Channel",
            labels=metrics.channels,
            datasets=[{"label": "ROAS", "data": metrics.roas}],
            y_label="ROAS (x)"
        )
    ))

    # 3. Spend vs Revenue grouped bar
    slides.append(SlideContent(
        title="Spend vs. Revenue by Channel",
        layout="chart",
        bullets=["Revenue consistently outpaces spend across all channels"],
        chart_data=ChartData(
            chart_type="grouped_bar",
            title="Spend vs. Revenue by Channel (£)",
            labels=metrics.channels,
            datasets=[
                {"label": "Spend",   "data": metrics.spend},
                {"label": "Revenue", "data": metrics.revenue},
            ],
            y_label="Value (£)"
        )
    ))

    # 4. Spend share (pie)
    slides.append(SlideContent(
        title="Spend Share by Channel",
        layout="chart",
        bullets=["Sponsored Products accounts for the majority of investment"],
        chart_data=ChartData(
            chart_type="pie",
            title="Spend Share by Channel",
            labels=metrics.channels,
            datasets=[{"data": metrics.spend}]
        )
    ))

    # 5. Trend – Spend & Revenue over 6 months (line)
    slides.append(SlideContent(
        title="6-Month Performance Trend",
        layout="chart",
        bullets=["Strong revenue growth driven by increasing spend efficiency"],
        chart_data=ChartData(
            chart_type="line",
            title="6-Month Spend & Revenue Trend (£)",
            labels=metrics.trend_months,
            datasets=[
                {"label": "Spend",   "data": metrics.trend_spend},
                {"label": "Revenue", "data": metrics.trend_revenue},
            ],
            y_label="Value (£)"
        )
    ))

    return slides


# ─── Claude prompt → slide content ───────────────────────────────────────────

def generate_slides_from_prompt(
    prompt: str,
    company: str = None,
    context_data: Optional[dict] = None
) -> List[dict]:
    if not ANTHROPIC_KEY:
        return []

    client = anthropic.Anthropic(api_key=ANTHROPIC_KEY)

    system = """Generate presentation slides as JSON. Output ONLY valid JSON array.

Format:
[
  {
    "title": "Slide Title",
    "bullets": ["Point 1", "Point 2", "Point 3"],
    "notes": "Speaker notes",
    "layout": "content"
  },
  ...
]

Layout options: "content" | "section" | "closing"

Rules:
- 4-6 slides max
- 3-5 bullets per slide
- Keep bullets concise (under 15 words)
- Professional business tone
- Use "section" layout for divider/transition slides (1-2 bullets max)
- Use "closing" layout for the final slide
"""

    user_prompt = prompt
    if company:
        user_prompt = f"For {company}: {prompt}"
    if context_data:
        user_prompt += f"\n\nContext data: {json.dumps(context_data)}"

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=2000,
            system=system,
            messages=[{"role": "user", "content": user_prompt}]
        )
        text = response.content[0].text
        if "```json" in text:
            text = text.split("```json")[1].split("```")[0]
        elif "```" in text:
            text = text.split("```")[1].split("```")[0]
        return json.loads(text)
    except Exception as e:
        print(f"Claude error: {e}")
        return []


# ─── Slide builders ───────────────────────────────────────────────────────────

def build_title_slide(prs: "Presentation", request: PPTRequest):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    W = prs.slide_width
    H = prs.slide_height

    # Gradient background: deep navy → mid-blue
    add_gradient_background(slide, (0x00, 0x1A, 0x33), (0x0A, 0x3A, 0x6E), angle=135)

    # TAU teal left accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAU_ACCENT
    bar.line.fill.background()

    # Company name (small caps style)
    if request.company:
        add_text_box(slide, request.company.upper(),
                     Inches(0.4), Inches(1.2), Inches(12), Inches(0.5),
                     size=13, bold=True, color=TAU_ACCENT, align=PP_ALIGN.LEFT)

    # Main title
    add_text_box(slide, request.title,
                 Inches(0.4), Inches(2.0), Inches(12), Inches(1.8),
                 size=48, bold=True, color=TAU_WHITE, align=PP_ALIGN.LEFT)

    # Accent rule under title
    add_accent_line(slide, Inches(3.9), width_inches=11.5, color=TAU_ACCENT, thickness_pt=2)

    # Subtitle
    if request.subtitle:
        add_text_box(slide, request.subtitle,
                     Inches(0.4), Inches(4.1), Inches(10), Inches(1),
                     size=22, color=TAU_LIGHT, align=PP_ALIGN.LEFT)

    # Date bottom-right
    date_str = datetime.now().strftime("%B %Y")
    add_text_box(slide, date_str,
                 Inches(9), Inches(6.5), Inches(3.8), Inches(0.5),
                 size=11, color=TAU_LIGHT, align=PP_ALIGN.RIGHT, italic=True)

    # Footer bar
    add_bottom_bar(slide, H)


def build_content_slide(prs: "Presentation", slide_data: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    # Dark background with subtle lighter band at top
    add_gradient_background(slide, (0x05, 0x14, 0x28), (0x0A, 0x28, 0x48), angle=160)

    # Top accent bar (thin)
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TAU_ACCENT
    top_bar.line.fill.background()

    # Title
    add_text_box(slide, slide_data.title,
                 Inches(0.5), Inches(0.25), Inches(12), Inches(0.85),
                 size=32, bold=True, color=TAU_WHITE, align=PP_ALIGN.LEFT)

    # Accent divider under title
    add_accent_line(slide, Inches(1.15), width_inches=12, color=TAU_BLUE, thickness_pt=1.5)

    # Bullets
    if slide_data.bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(11.8), Inches(5.4))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        bullet_colors = [TAU_WHITE, TAU_LIGHT, TAU_LIGHT, TAU_LIGHT, TAU_LIGHT]
        for i, bullet in enumerate(slide_data.bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(21)
            p.font.color.rgb = bullet_colors[min(i, len(bullet_colors) - 1)]
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0
            # Use a bullet character
            pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
            buChar_xml = '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="▸"/>'
            try:
                buChar = parse_xml(buChar_xml)
                pPr.append(buChar)
            except Exception:
                p.text = f"▸  {bullet}"

    add_bottom_bar(slide, H)


def build_chart_slide(prs: "Presentation", slide_data: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    add_gradient_background(slide, (0x05, 0x14, 0x28), (0x0D, 0x1B, 0x2A), angle=150)

    # Top accent bar
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TAU_ACCENT
    top_bar.line.fill.background()

    # Slide title
    add_text_box(slide, slide_data.title,
                 Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                 size=28, bold=True, color=TAU_WHITE, align=PP_ALIGN.LEFT)

    # Subtitle / insight text
    if slide_data.bullets:
        insight = slide_data.bullets[0]
        add_text_box(slide, f"↗  {insight}",
                     Inches(0.5), Inches(0.85), Inches(12), Inches(0.5),
                     size=14, color=TAU_ACCENT, align=PP_ALIGN.LEFT, italic=True)

    # Chart
    if slide_data.chart_data and MATPLOTLIB_AVAILABLE:
        img_bytes = chart_bytes_from_spec(slide_data.chart_data)
        if img_bytes:
            # Chart area: x=0.4, y=1.4, w=12.5, h=5.5
            embed_chart_image(slide, img_bytes,
                              Inches(0.4), Inches(1.4),
                              Inches(12.5), Inches(5.5))

    add_bottom_bar(slide, H)


def build_section_slide(prs: "Presentation", slide_data: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    add_gradient_background(slide, (0x00, 0x2B, 0x55), (0x00, 0x4A, 0x8A), angle=120)

    # Vertical left accent
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TAU_ACCENT
    bar.line.fill.background()

    add_text_box(slide, slide_data.title,
                 Inches(0.5), Inches(2.8), Inches(12), Inches(1.2),
                 size=40, bold=True, color=TAU_WHITE, align=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(4.1), color=TAU_ACCENT, thickness_pt=2)

    if slide_data.bullets:
        add_text_box(slide, slide_data.bullets[0],
                     Inches(1), Inches(4.35), Inches(11), Inches(0.7),
                     size=18, color=TAU_LIGHT, align=PP_ALIGN.CENTER, italic=True)

    add_bottom_bar(slide, H)


def build_closing_slide(prs: "Presentation", slide_data: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    add_gradient_background(slide, (0x00, 0x1A, 0x33), (0x00, 0x33, 0x66), angle=135)

    # Gold accent top bar
    top_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), W, Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TAU_GOLD
    top_bar.line.fill.background()

    add_text_box(slide, slide_data.title,
                 Inches(0.5), Inches(2.2), Inches(12), Inches(1.5),
                 size=44, bold=True, color=TAU_WHITE, align=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(3.8), color=TAU_GOLD, thickness_pt=2)

    if slide_data.bullets:
        body = "\n".join(slide_data.bullets)
        add_text_box(slide, body,
                     Inches(1.5), Inches(4.0), Inches(10), Inches(2.0),
                     size=18, color=TAU_LIGHT, align=PP_ALIGN.CENTER)

    add_bottom_bar(slide, H, label="TAU — The Independent Marketing Intelligence Architect  |  tau.ai")


# ─── Main PPT builder ─────────────────────────────────────────────────────────

def create_ppt(request: PPTRequest) -> io.BytesIO:
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=500, detail="python-pptx not installed")

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title slide ──
    build_title_slide(prs, request)

    # ── Content / chart slides ──
    slides_content: List[SlideContent] = list(request.slides or [])

    # Auto-generate from retail metrics
    if request.retail_metrics and MATPLOTLIB_AVAILABLE:
        chart_slides = build_retail_chart_slides(request.retail_metrics)
        slides_content = chart_slides + slides_content

    # Generate textual slides from Claude prompt if no slides provided
    if request.prompt and not slides_content:
        generated = generate_slides_from_prompt(request.prompt, request.company, request.context_data)
        slides_content = [SlideContent(**s) for s in generated]

    for s in slides_content:
        layout = s.layout
        if layout == "chart":
            build_chart_slide(prs, s)
        elif layout == "section":
            build_section_slide(prs, s)
        elif layout == "closing":
            build_closing_slide(prs, s)
        else:
            build_content_slide(prs, s)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# ─── Routes ───────────────────────────────────────────────────────────────────

@router.post("/ppt/generate")
async def generate_ppt(request: PPTRequest):
    """Generate a professional TAU-branded PowerPoint presentation."""
    try:
        buffer = create_ppt(request)
        filename = f"{request.title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/ppt/preview")
async def preview_ppt(request: PPTRequest):
    """Preview slide content without generating file."""
    slides_content = list(request.slides or [])

    if request.retail_metrics:
        chart_slides = build_retail_chart_slides(request.retail_metrics)
        slides_content = chart_slides + slides_content

    if request.prompt and not slides_content:
        generated = generate_slides_from_prompt(request.prompt, request.company, request.context_data)
        slides_content = generated

    return {
        "title": request.title,
        "subtitle": request.subtitle,
        "slides": slides_content,
        "slide_count": len(slides_content) + 1,
        "matplotlib_available": MATPLOTLIB_AVAILABLE,
        "pptx_available": PPTX_AVAILABLE,
    }


@router.post("/ppt/demo")
async def generate_demo_ppt():
    """
    Generate a demo retail media performance deck with sample metrics.
    Tests chart integration, gradient backgrounds, and all layout types.
    """
    request = PPTRequest(
        title="Retail Media Performance Review",
        subtitle="Q1 2026 — Sponsored & Display Campaign Analysis",
        company="Nectar360",
        retail_metrics=RetailMetrics(),  # default sample data
        slides=[
            SlideContent(
                title="Executive Summary",
                layout="content",
                bullets=[
                    "Total retail media investment: £220k in March 2026",
                    "Blended ROAS of 5.8x across all channels — up 12% MoM",
                    "Sponsored Products remains top-performing channel at 7.0x ROAS",
                    "Video CPM down 18% while impressions grew 23%",
                    "Incremental revenue attributed to retail media: £1.54M",
                ],
                notes="Lead with the headline numbers. Audience expects ROAS focus."
            ),
            SlideContent(
                title="Strategic Recommendations",
                layout="section",
                bullets=["Shift budget toward highest-ROAS channels and scale seasonal peaks"],
            ),
            SlideContent(
                title="Recommended Actions",
                layout="content",
                bullets=[
                    "Increase Sponsored Products budget by 15% in April",
                    "Expand Display retargeting to mid-funnel audiences",
                    "Test Video creative A/B for improved VTR",
                    "Implement dayparting to reduce off-peak CPC waste",
                    "Launch new category page placements ahead of summer",
                ],
            ),
            SlideContent(
                title="Next Steps & Contact",
                layout="closing",
                bullets=[
                    "April planning session — w/c 7 April",
                    "Budget reallocation sign-off by 31 March",
                    "Contact: strategy@tau.ai  |  tau.ai",
                ],
            ),
        ]
    )
    try:
        buffer = create_ppt(request)
        filename = f"retail_media_demo_{datetime.now().strftime('%Y%m%d')}.pptx"
        return StreamingResponse(
            buffer,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
