#!/usr/bin/env python3
"""Build Agent Builder Comparison PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_TEXT = RGBColor(0x1a, 0x1a, 0x1a)
HEADER_BG = RGBColor(0x2d, 0x2d, 0x4e)
ROW_ALT = RGBColor(0xF8, 0xF8, 0xFC)
GREEN = RGBColor(0x00, 0xAA, 0x55)
RED = RGBColor(0xCC, 0x33, 0x33)
ORANGE = RGBColor(0xFF, 0x99, 0x00)

SCREENSHOTS = "/home/r2/nectar-retail-media/screenshots"


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=14, color=WHITE, bold=False, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    return p


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

add_textbox(slide, 1, 1.5, 11, 1.5,
    "Build Your Own Agent", 48, ACCENT, True, PP_ALIGN.CENTER)
tf = add_textbox(slide, 1, 3, 11, 1,
    "Platform Comparison & Recommendation", 28, WHITE, False, PP_ALIGN.CENTER)
add_para(tf, "", 14, WHITE)
add_para(tf, "NinjaCat  |  Akkio  |  MindStudio  |  Relevance AI  |  Lyzr", 18, MID_GRAY, False)
add_para(tf, "Google Vertex AI  |  Google Workspace Studio  |  Microsoft Copilot Studio  |  Claude", 18, MID_GRAY, False)
add_para(tf, "", 14, WHITE)
add_para(tf, "TAU  |  March 2026", 16, ACCENT, False)


# ============================================================
# SLIDE 2: Agent Types Framework
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, 0.5, 0.3, 12, 0.8,
    "Agent Types Framework", 36, ACCENT, True)
add_textbox(slide, 0.5, 1.0, 12, 0.5,
    "How we categorise agents by function", 16, MID_GRAY)

agent_types = [
    ("Data Collectors", "Gather, ingest & normalise data from APIs, feeds, platforms",
     "Campaign metrics, CRM sync, competitor monitoring, web scraping", ACCENT),
    ("Alerting / Monitoring", "Watch for thresholds, anomalies, changes & notify",
     "Budget pacing alerts, performance drops, brand mention detection", RGBColor(0xFF, 0x66, 0x66)),
    ("Analysis / Insights", "Process data to surface patterns, trends & recommendations",
     "Cross-channel attribution, audience insights, creative performance", RGBColor(0x66, 0xBB, 0xFF)),
    ("Action / Optimisation", "Take autonomous actions based on analysis",
     "Pause underperformers, reallocate budget, adjust bids, A/B test", RGBColor(0xFF, 0xCC, 0x00)),
    ("Content / Creative", "Generate reports, copy, emails, presentations",
     "Client reports, ad copy, email drafts, social posts", RGBColor(0xCC, 0x66, 0xFF)),
    ("Orchestration", "Coordinate multiple agents in workflows",
     "Multi-step pipelines: collect > analyse > alert > act > report", RGBColor(0x66, 0xFF, 0x99)),
]

for i, (name, desc, examples, color) in enumerate(agent_types):
    y = 1.6 + i * 0.9
    # Color bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.15), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Name
    add_textbox(slide, 0.8, y - 0.05, 3, 0.4, name, 18, color, True)
    # Desc
    add_textbox(slide, 0.8, y + 0.3, 5.5, 0.4, desc, 12, WHITE)
    # Examples
    add_textbox(slide, 6.8, y, 6, 0.7, examples, 12, MID_GRAY)


# ============================================================
# SLIDE 3: Comparison Matrix
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, 0.3, 0.2, 12, 0.7,
    "Platform Comparison Matrix", 32, DARK_TEXT, True)

headers = ["Platform", "Data\nCollectors", "Alerts /\nMonitoring", "Analysis /\nInsights", "Action /\nOptimise", "Content /\nCreative", "Orchestration", "No-Code", "White\nLabel", "Free\nTier"]

# Y = yes, P = partial, N = no
rows = [
    ["NinjaCat",       "Y","Y","Y","Y","Y","P", "Y","N","N"],
    ["Akkio",          "Y","P","Y","N","Y","N", "Y","Y","Trial"],
    ["MindStudio",     "Y","Y","Y","P","Y","P", "Y","N","Y"],
    ["Relevance AI",   "Y","P","Y","P","Y","Y", "Low","N","Y"],
    ["Lyzr",           "Y","P","Y","P","Y","Y", "Y","N","?"],
    ["Google Vertex",  "Y","Y","Y","Y","Y","Y", "Low","N","GCP"],
    ["Google Workspace","P","P","Y","P","Y","N", "Y","N","Plan"],
    ["Copilot Studio", "Y","Y","Y","Y","Y","Y", "Y","N","Trial"],
    ["Claude",         "Y","P","Y","Y","Y","Y", "Code","N","Plan"],
]

col_widths = [1.8, 1.0, 1.0, 1.0, 1.0, 1.0, 1.1, 0.8, 0.8, 0.8]
start_x = 0.4
start_y = 1.0
row_h = 0.5
header_h = 0.6

# Header row
x = start_x
for j, (h, w) in enumerate(zip(headers, col_widths)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(start_y), Inches(w), Inches(header_h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = h
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    x += w

# Data rows
for i, row in enumerate(rows):
    x = start_x
    y = start_y + header_h + i * row_h
    bg_color = ROW_ALT if i % 2 == 0 else WHITE
    for j, (val, w) in enumerate(zip(row, col_widths)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(row_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if val == "Y":
            p.text = "Yes"
            p.font.color.rgb = GREEN
            p.font.bold = True
        elif val == "N":
            p.text = "No"
            p.font.color.rgb = RED
        elif val == "P":
            p.text = "Partial"
            p.font.color.rgb = ORANGE
        else:
            p.text = val
            p.font.color.rgb = DARK_TEXT
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += w

# Legend
add_textbox(slide, 0.4, 6.5, 10, 0.5,
    "Y = Full support  |  P = Partial/limited  |  N = Not available  |  Low = Low-code (some coding needed)  |  Code = Code-first",
    10, MID_GRAY)


# ============================================================
# SLIDE 4: Limitations Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
    "Key Limitations by Platform", 32, ACCENT, True)

limitations = [
    ("NinjaCat", "Marketing only - no general-purpose agents. Closed ecosystem. No white-label. Premium pricing, no free tier."),
    ("Akkio", "Analytics/BI focused - no action agents (can't pause campaigns, adjust bids). No alerting. No multi-agent orchestration."),
    ("MindStudio", "No native data connectors for ad platforms. No white-label. Orchestration is basic (no true multi-agent teams)."),
    ("Relevance AI", "Sales/GTM focus - not built for media/advertising. Credit-based pricing scales fast. 20% markup on LLM costs."),
    ("Lyzr", "Enterprise focus - overkill for small teams. Limited pre-built marketing templates. Pricing not transparent."),
    ("Google Vertex", "Requires GCP expertise. Steep learning curve. Code-heavy for advanced agents. Costs can spiral with scale."),
    ("Google Workspace", "Limited to Workspace data. No ad platform integrations. No action agents. Basic agent capabilities."),
    ("Copilot Studio", "Microsoft ecosystem lock-in. Limited ad platform connectors. Enterprise pricing. Complex for simple use cases."),
    ("Claude", "Code-first - not for non-technical users. No native ad platform connectors. No built-in scheduling/triggers."),
]

for i, (name, lim) in enumerate(limitations):
    y = 1.1 + i * 0.68
    add_textbox(slide, 0.5, y, 2.2, 0.5, name, 13, ACCENT, True)
    add_textbox(slide, 2.8, y, 10, 0.6, lim, 11, WHITE)


# ============================================================
# SLIDES 5-13: Individual Platform Slides with Screenshots
# ============================================================

platforms = [
    {
        "name": "NinjaCat",
        "screenshot": "ninjacat-agents.png",
        "tagline": "AI Agents Built for Marketing",
        "url": "ninjacat.io",
        "agents": [
            "Data Collectors: Auto-ingest from Google Ads, Meta, GA4, LinkedIn, CRM, 100+ sources",
            "Alerts: Spend pacing, CPA targets, impression share, engagement anomalies",
            "Analysis: Cross-account benchmarking, KPI trend digests, creative performance",
            "Actions: Pause underperformers, reallocate budgets, adjust bids automatically",
            "Content: Auto-generate dashboards, performance summaries, client briefs",
        ],
        "strengths": "Most complete marketing agent platform. Data Cloud unifies all sources. Agents are pre-trained on ad-tech patterns.",
        "weaknesses": "Marketing only. No general-purpose agents. Closed ecosystem. No free tier.",
        "verdict": "BEST FOR: Agencies wanting turnkey marketing AI agents",
    },
    {
        "name": "Akkio",
        "screenshot": "akkio-platform.png",
        "tagline": "White-Label Generative BI for Agencies",
        "url": "akkio.com",
        "agents": [
            "Data Collectors: Connect client data sources, auto-clean and prepare",
            "Analysis: Chat Explore (natural language queries), ML predictions, forecasts",
            "Content: Branded dashboards, insights reports, visualisations",
            "Alerts: Limited - no proactive monitoring or threshold alerts",
            "Actions: None - analytics only, no campaign actions",
        ],
        "strengths": "Only true white-label option. Agencies embed under their own brand. Great for client self-service analytics.",
        "weaknesses": "No action agents. No alerting. No orchestration. Analytics/BI only - not a full agent platform.",
        "verdict": "BEST FOR: Agencies wanting to resell AI analytics under their brand",
    },
    {
        "name": "MindStudio",
        "screenshot": "mindstudio.png",
        "tagline": "200+ AI Models, No-Code Agent Builder",
        "url": "mindstudio.ai",
        "agents": [
            "Data Collectors: API calls, web scraping, SQL queries, webhook triggers",
            "Alerts: Scheduled agents + webhook triggers for monitoring workflows",
            "Analysis: Data viz agents, research agents, market analysis",
            "Content: Email drafts, reports, social posts, ad copy — 100+ templates",
            "Actions: Partial - can call APIs but no native ad platform connectors",
        ],
        "strengths": "Most model-agnostic (200+ models). Dynamic Tool Use. Deploy anywhere (web, email, API, browser extension).",
        "weaknesses": "No native ad platform integrations. No white-label. Basic orchestration only.",
        "verdict": "BEST FOR: Teams wanting maximum flexibility across AI models and use cases",
    },
    {
        "name": "Relevance AI",
        "screenshot": "relevance-ai.png",
        "tagline": "Multi-Agent AI Workforce for Sales & GTM",
        "url": "relevanceai.com",
        "agents": [
            "Data Collectors: Lead research, CRM enrichment, competitor analysis, web research",
            "Alerts: Limited - notification on task completion, no real-time monitoring",
            "Analysis: Deal intelligence, campaign performance, market data analysis",
            "Content: Proposals, email personalisation, content repurposing, blog drafts",
            "Orchestration: Multi-agent teams (research > verify > write > distribute)",
        ],
        "strengths": "Best multi-agent orchestration. 100+ templates. Agents collaborate as teams on complex tasks.",
        "weaknesses": "Sales/GTM focus. Credit pricing scales fast. 20% LLM markup without own API keys.",
        "verdict": "BEST FOR: Sales teams wanting autonomous BDR and research agents",
    },
    {
        "name": "Lyzr",
        "screenshot": "lyzr-agent-studio.png",
        "tagline": "Enterprise AI Agent Studio with Safe AI",
        "url": "lyzr.ai",
        "agents": [
            "Data Collectors: API connectors, custom tool integration, data pipelines",
            "Alerts: Partial - workflow-triggered notifications",
            "Analysis: LLM + ML hybrid for better accuracy, blueprints for support/ops/finance",
            "Content: Report generation, customer communications",
            "Orchestration: Solo agents, multi-agent systems, full orchestration",
        ],
        "strengths": "Enterprise compliance (SOC2, GDPR, HIPAA, ISO 27001). Safe AI built into core. Flexible deployment.",
        "weaknesses": "Enterprise-focused - overkill for small teams. Limited marketing templates. Opaque pricing.",
        "verdict": "BEST FOR: Regulated industries needing compliance-first agent platform",
    },
    {
        "name": "Google Vertex AI Agent Builder",
        "screenshot": "google-vertex-agent-builder.png",
        "tagline": "Enterprise-Grade Agent Development on GCP",
        "url": "cloud.google.com/products/agent-builder",
        "agents": [
            "Data Collectors: BigQuery, Cloud Storage, 100+ connectors, Dataproc for large-scale",
            "Alerts: Cloud Monitoring, auto-monitoring of agent behaviour and reasoning",
            "Analysis: Multi-agent financial research, RAG agents, custom ML pipelines",
            "Actions: Enterprise system connectors (ERP, procurement, HR), workflow automation",
            "Orchestration: Agent Engine (managed runtime), directed graph workflows, ADK framework",
        ],
        "strengths": "Most powerful and scalable. Open-source ADK. Full lifecycle management. Global scale.",
        "weaknesses": "Requires GCP expertise. Steep learning curve. Code-heavy. Costs can spiral.",
        "verdict": "BEST FOR: Engineering teams building production-grade agent systems at scale",
    },
    {
        "name": "Google Workspace Studio",
        "screenshot": "google-workspace-studio.png",
        "tagline": "No-Code Agents for Google Workspace",
        "url": "workspace.google.com/studio",
        "agents": [
            "Data Collectors: Partial - Gmail, Drive, Sheets, Calendar data only",
            "Alerts: Partial - can trigger on Workspace events",
            "Analysis: Summarise docs, analyse spreadsheets, extract email insights",
            "Content: Draft emails, create docs, generate presentations",
            "Actions: Partial - Workspace actions only (send email, create doc, update sheet)",
        ],
        "strengths": "Zero friction for Google Workspace orgs. Build in minutes. No coding. Free with Business/Enterprise.",
        "weaknesses": "Limited to Workspace data. No ad platforms. No real agent orchestration. Basic capabilities.",
        "verdict": "BEST FOR: Quick internal automations within Google Workspace",
    },
    {
        "name": "Microsoft Copilot Studio",
        "screenshot": "copilot-studio.png",
        "tagline": "Create, Customise and Launch AI Agents",
        "url": "microsoft.com/copilot-studio",
        "agents": [
            "Data Collectors: SharePoint, OneDrive, Planner, Forms, Dynamics 365, custom connectors",
            "Alerts: Event triggers (file created, task completed, form submitted, schedule-based)",
            "Analysis: Conversational analytics over M365 data, Power BI integration",
            "Actions: Autonomous agents that monitor, decide, and execute without user prompts",
            "Orchestration: Agent flows, multi-step automations, Power Automate integration",
        ],
        "strengths": "Best autonomous agent triggers. Deep M365 integration. No-code. Enterprise-grade audit logging.",
        "weaknesses": "Microsoft ecosystem lock-in. Limited ad platform connectors. Enterprise pricing.",
        "verdict": "BEST FOR: Microsoft-heavy organisations wanting autonomous background agents",
    },
    {
        "name": "Claude (Anthropic)",
        "screenshot": "claude-cowork.png",
        "tagline": "AI Agents via Cowork & Agent SDK",
        "url": "claude.com/solutions/agents",
        "agents": [
            "Data Collectors: Web scraping, API calls, file processing, MCP server integrations",
            "Alerts: Partial - no built-in scheduling, needs external triggers",
            "Analysis: Most capable reasoning — multi-step analysis, code execution, research",
            "Actions: File editing, code execution, tool calling, system commands",
            "Orchestration: Agent Teams (up to 10 parallel sub-agents), Agent SDK, Skills marketplace",
        ],
        "strengths": "Most capable reasoning model. Best for complex multi-step tasks. Open SDK. Agent Teams for parallelism.",
        "weaknesses": "Code-first — not for non-technical users. No native ad connectors. No scheduling/triggers.",
        "verdict": "BEST FOR: Developer teams building custom intelligent workflows",
    },
]

for p in platforms:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)

    # Header
    add_textbox(slide, 0.5, 0.2, 8, 0.6, p["name"], 36, ACCENT, True)
    add_textbox(slide, 0.5, 0.75, 8, 0.4, p["tagline"], 16, MID_GRAY)
    add_textbox(slide, 0.5, 1.05, 8, 0.3, p["url"], 12, RGBColor(0x88, 0x88, 0xBB))

    # Screenshot
    img_path = os.path.join(SCREENSHOTS, p["screenshot"])
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.5), Inches(0.3), Inches(5.5))

    # Agent types
    tf = add_textbox(slide, 0.5, 1.5, 6.8, 4, "Agent Capabilities:", 14, WHITE, True)
    for agent_line in p["agents"]:
        add_para(tf, agent_line, 11, WHITE, False, Pt(6))

    # Strengths
    add_para(tf, "", 8, WHITE)
    add_para(tf, "Strengths:", 13, GREEN, True, Pt(10))
    add_para(tf, p["strengths"], 11, WHITE, False, Pt(4))

    # Weaknesses
    add_para(tf, "Limitations:", 13, RGBColor(0xFF, 0x66, 0x66), True, Pt(10))
    add_para(tf, p["weaknesses"], 11, WHITE, False, Pt(4))

    # Verdict bar
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.5), Inches(12), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = HEADER_BG
    shape.line.fill.background()
    tf2 = shape.text_frame
    tf2.paragraphs[0].text = p["verdict"]
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.color.rgb = ACCENT
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE


# ============================================================
# SLIDE 14: Recommendation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_textbox(slide, 0.5, 0.3, 12, 0.7,
    "Recommendation", 36, ACCENT, True)

tf = add_textbox(slide, 0.5, 1.2, 12, 5.5, "", 14, WHITE)

add_para(tf, "For a Retail Media / Ad-Tech Context:", 22, ACCENT, True, Pt(0))
add_para(tf, "", 10, WHITE)

add_para(tf, "Tier 1: Purpose-Built for Marketing", 18, GREEN, True, Pt(10))
add_para(tf, "NinjaCat — if you want the most complete out-of-the-box marketing agent platform (data collection, alerts, analysis, actions, reporting all built in). Best if you need agents that understand ad-tech natively.", 13, WHITE, False, Pt(6))

add_para(tf, "", 8, WHITE)
add_para(tf, "Tier 2: White-Label & Client-Facing", 18, RGBColor(0xFF, 0xCC, 0x00), True, Pt(10))
add_para(tf, "Akkio — if the goal is giving clients self-service analytics under your brand. Not a full agent platform but the only real white-label option.", 13, WHITE, False, Pt(6))

add_para(tf, "", 8, WHITE)
add_para(tf, "Tier 3: Build Custom (Most Flexible)", 18, RGBColor(0x66, 0xBB, 0xFF), True, Pt(10))
add_para(tf, "Google Vertex AI — if you have engineering resource and want full control at scale.", 13, WHITE, False, Pt(6))
add_para(tf, "Copilot Studio — if you're a Microsoft shop wanting autonomous background agents with triggers.", 13, WHITE, False, Pt(6))
add_para(tf, "Claude Agent SDK — if you want the most capable reasoning for complex analysis tasks.", 13, WHITE, False, Pt(6))

add_para(tf, "", 8, WHITE)
add_para(tf, "Tier 4: General Purpose / Niche", 18, MID_GRAY, True, Pt(10))
add_para(tf, "MindStudio (model flexibility), Relevance AI (sales teams), Lyzr (regulated industries), Google Workspace Studio (internal automations)", 13, WHITE, False, Pt(6))

# Bottom note
add_textbox(slide, 0.5, 6.8, 12, 0.5,
    "Key insight: No single platform covers all 6 agent types well for ad-tech. NinjaCat is closest but lacks white-label. A hybrid approach (NinjaCat + Akkio) may be optimal.",
    12, ACCENT, True)


# ============================================================
# Save
# ============================================================
out_path = "/home/r2/nectar-retail-media/agent-builder-comparison.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
